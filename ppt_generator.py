from pptx import Presentation
from pptx.util import Pt, Inches
from datetime import datetime
import re

TITLE_FONT_SIZE = 32
SUBTITLE_FONT_SIZE = 24
BODY_FONT_SIZE = 18

TITLE_SLIDE_LAYOUT = 0
SECTION_HEADER_LAYOUT = 2
TITLE_AND_CONTENT_LAYOUT = 5

def add_bullet_point(paragraph, text, level=0, font_size=BODY_FONT_SIZE, bold=False):
    """Add a bullet point with consistent formatting."""
    paragraph.text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.level = level
    return paragraph

def add_section_header(slide, title, subtitle=None):
    """Add title (and optional subtitle) to a slide."""
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(TITLE_FONT_SIZE)
    title_shape.text_frame.paragraphs[0].font.bold = True
    if subtitle:
        subtitle_shape = slide.shapes.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(SUBTITLE_FONT_SIZE)

def create_content_box(slide, left, top, width, height):
    """Return a text-frame inside a new textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf

def format_financial_metric(value):
    return value if value else None

def format_percentage(value):
    return value if value else None

def generate_ppt(document_data, output_path):
    prs = Presentation()

    analysis = document_data.get('analysis', {})
    company_info = analysis.get('company_info', {})
    financial_metrics = document_data.get('financial_metrics', {})
    market_info = analysis.get('market_info', {})

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    company_name = company_info.get('name', 'Investment Opportunity')
    title_shape.text = company_name
    subtitle_shape.text = f"Investment Overview\n{datetime.now().strftime('%B %Y')}"
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title_shape = slide.shapes.title
    title_shape.text = "Company Overview"
    title_shape.text_frame.paragraphs[0].font.size = Pt(TITLE_FONT_SIZE)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    content_tf = create_content_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    
    content_tf.add_paragraph()
    
    p = content_tf.add_paragraph()
    p.text = "Company Overview"; p.font.size = Pt(SUBTITLE_FONT_SIZE); p.font.bold = True
    if company_info.get('location'):
        add_bullet_point(content_tf.add_paragraph(), f"Headquarters: {company_info['location']}")
    if market_info.get('market_size'):
        add_bullet_point(content_tf.add_paragraph(), f"Market Size: {market_info['market_size']}")
    if market_info.get('growth_rate'):
        add_bullet_point(content_tf.add_paragraph(), f"Market Growth: {market_info['growth_rate']}")
    
    content_tf.add_paragraph()
    
    p = content_tf.add_paragraph()
    p.text = "Business Model"; p.font.size = Pt(SUBTITLE_FONT_SIZE); p.font.bold = True
    if company_info.get('business_model'):
        p = content_tf.add_paragraph()
        p.text = company_info['business_model']
        p.font.size = Pt(BODY_FONT_SIZE)

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    add_section_header(slide, "Key Financials")
    content_tf = create_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4))

    metrics = {
        "Revenue": format_financial_metric(financial_metrics.get('revenue')),
        "EBITDA": format_financial_metric(financial_metrics.get('ebitda')),
        "Revenue Growth": format_percentage(financial_metrics.get('revenue_growth')),
        "EBITDA Margin": format_percentage(financial_metrics.get('ebitda_margin'))
    }
    displayed = False
    for label, value in metrics.items():
        if value:
            add_bullet_point(content_tf.add_paragraph(), f"{label}: {value}", bold=False)
            displayed = True
    if not displayed:
        p = content_tf.add_paragraph()
        p.text = "Financial information not available in the document"
        p.font.italic = True
        p.font.size = Pt(BODY_FONT_SIZE)

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    add_section_header(slide, "Investment Rationale")
    content_tf = create_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4))

    highlights = analysis.get('key_highlights', [])
    fallback_points = []
    if not highlights:
        if company_info.get('business_model'):
            fallback_points.append(company_info['business_model'])
        if market_info.get('market_size'):
            fallback_points.append(f"Strong market presence in a {market_info['market_size']} industry")
        if market_info.get('growth_rate'):
            fallback_points.append(f"Operating in a high-growth market ({market_info['growth_rate']} growth)")
        if market_info.get('competition'):
            fallback_points.append(market_info['competition'])
        if financial_metrics.get('revenue_growth'):
            fallback_points.append(f"Demonstrated strong growth with {financial_metrics['revenue_growth']} revenue increase")
        highlights = fallback_points

    for idx, point in enumerate(highlights[:5], 1):
        add_bullet_point(content_tf.add_paragraph(), f"{idx}. {point}", bold=False)

    if market_info.get('competition') or market_info.get('market_size'):
        slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
        add_section_header(slide, "Market Position")
        content_tf = create_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4))

        if market_info.get('market_size'):
            p = content_tf.add_paragraph()
            p.text = "Market Opportunity"; p.font.size = Pt(SUBTITLE_FONT_SIZE); p.font.bold = True
            add_bullet_point(content_tf.add_paragraph(), f"Total Addressable Market: {market_info['market_size']}")
            if market_info.get('growth_rate'):
                add_bullet_point(content_tf.add_paragraph(), f"Market Growth Rate: {market_info['growth_rate']}")

        if market_info.get('competition'):
            content_tf.add_paragraph()
            p = content_tf.add_paragraph()
            p.text = "Competitive Position"; p.font.size = Pt(SUBTITLE_FONT_SIZE); p.font.bold = True
            add_bullet_point(content_tf.add_paragraph(), market_info['competition'])

    prs.save(output_path)